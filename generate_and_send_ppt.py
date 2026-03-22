#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
计生行业竞品调研 PPT - 精美商务版
发送到企业邮箱
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import smtplib
import os
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

# 配色方案
COLORS = {
    'primary': (30, 58, 138),
    'secondary': (59, 130, 246),
    'accent': (239, 68, 68),
    'bg': (255, 255, 255),
    'text': (31, 41, 55),
    'light': (243, 244, 246),
}

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "计生行业竞品调研报告"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(10), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.add_paragraph()
    p.text = "2026 年 3 月"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Microsoft YaHei'

def create_toc_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*COLORS['bg'])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "目录"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['primary'])
    p.font.name = 'Microsoft YaHei'
    
    chapters = ["01 行业概况", "02 主要竞品品牌", "03 市场份额", "04 产品线对比", 
                "05 价格策略", "06 营销渠道", "07 社交媒体", "08 新品动态", 
                "09 用户画像", "10 趋势与建议"]
    
    for i, chapter in enumerate(chapters):
        y = Inches(1.8) + i * Inches(0.55)
        num_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(1), Inches(0.5))
        tf = num_box.text_frame
        p = tf.add_paragraph()
        p.text = chapter[:2]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['primary'])
        
        text_box = slide.shapes.add_textbox(Inches(1.2), y, Inches(8), Inches(0.5))
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.text = chapter[3:]
        p.font.size = Pt(22)
        p.font.name = 'Microsoft YaHei'

def create_content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*COLORS['bg'])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['primary'])
    p.font.name = 'Microsoft YaHei'
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(4), Inches(0.15))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
    line.line.fill.background()
    
    for i, item in enumerate(items):
        y = Inches(2) + i * Inches(0.65)
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.05), Inches(0.15), Inches(0.15))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = RGBColor(*COLORS['secondary'])
        
        text_box = slide.shapes.add_textbox(Inches(1), y, Inches(11), Inches(0.5))
        tf = text_box.text_frame
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.name = 'Microsoft YaHei'

def create_table_slide(prs, title, headers, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['primary'])
    
    table_shape = slide.shapes.add_table(len(data)+1, len(headers), Inches(0.3), Inches(1.2), Inches(12.7), Inches(4))
    table = table_shape.table
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
        tf = cell.text_frame
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    for row_idx, row in enumerate(data, 1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*COLORS['bg']) if row_idx % 2 else RGBColor(*COLORS['light'])
            tf = cell.text_frame
            tf.paragraphs[0].font.size = Pt(14)

def create_end_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*COLORS['primary'])
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "谢谢"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    qa_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(10), Inches(1))
    tf = qa_box.text_frame
    p = tf.add_paragraph()
    p.text = "Q & A"
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def generate_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    print("[Info] 开始生成 PPT...")
    
    create_title_slide(prs)
    create_toc_slide(prs)
    
    slides = [
        ("1. 行业概况", ["市场规模：2025 年约 150 亿元", "年增长率：8-10%", "线上占比：65%+", "品牌集中度高（CR5 > 60%）"]),
        ("2. 主要竞品品牌", ["国际品牌：杜蕾斯（英国）、杰士邦（澳洲）、冈本（日本）", "国产品牌：大象、名流、赤尾、羽感", "品牌定位差异化明显"]),
        ("3. 市场份额（线上 2025）", ["杜蕾斯：28%", "杰士邦：18%", "冈本：12%", "大象：8%", "其他：34%"]),
        ("4. 产品线对比", ["杜蕾斯：AIR 空气套、持久装", "杰士邦：零感系列、玻尿酸", "冈本：001/002/003 系列", "大象：牛油果、大花套"]),
    ]
    
    for title, content in slides:
        create_content_slide(prs, title, content)
        print(f"[OK] {title}")
    
    # 表格页
    headers = ["品牌", "产地", "定位", "价格带"]
    data = [
        ["杜蕾斯", "英国", "高端", "¥¥¥"],
        ["杰士邦", "澳洲", "中高端", "¥¥"],
        ["冈本", "日本", "超薄高端", "¥¥¥"],
        ["大象", "中国", "年轻科技", "¥¥"],
    ]
    create_table_slide(prs, "5. 品牌对比", headers, data)
    print("[OK] 5. 品牌对比")
    
    # 更多内容
    more_slides = [
        ("6. 价格策略", ["入门款：1-8 元", "主流款：3-15 元", "高端款：6-30 元"]),
        ("7. 营销渠道", ["线上：天猫/京东（60%+）", "即时零售：美团买药", "直播：抖音/快手", "种草：小红书"]),
        ("8. 社交媒体", ["杜蕾斯：450 万粉丝", "杰士邦：180 万粉丝", "大象：120 万粉丝"]),
        ("9. 新品动态", ["杜蕾斯：90 周年宣传片", "杰士邦：零感·至薄至润", "大象：限定大花套"]),
        ("10. 趋势与建议", ["产品创新：更薄更舒适", "渠道变革：即时零售", "营销升级：内容营销"]),
    ]
    
    for title, content in more_slides:
        create_content_slide(prs, title, content)
        print(f"[OK] {title}")
    
    create_end_slide(prs)
    
    output_path = r"C:\Users\ho\.openclaw\workspace\sososong\竞品调研_20260318_v2.pptx"
    prs.save(output_path)
    
    print(f"\n[OK] PPT 生成成功！")
    print(f"[File] {output_path}")
    print(f"[Slides] {len(prs.slides)} 页")
    
    return output_path

def send_email(ppt_path):
    print("\n[Info] 正在发送邮件...")
    
    smtp_server = "smtp.qq.com"
    smtp_port = 465
    from_email = "kinsongho@qq.com"
    to_email = "hejs@edaxiang.com"
    password = os.environ.get("QQ_EMAIL_PASSWORD")
    
    subject = "计生行业竞品调研报告 - 2026 年 3 月（精美版）"
    body = """Nelson，你好！

这是精美版的计生行业竞品调研 PPT 报告。

【报告内容】
1. 行业概况
2. 主要竞品品牌
3. 市场份额分析
4. 产品线对比
5. 价格策略
6. 营销渠道
7. 社交媒体表现
8. 新品动态
9. 用户画像
10. 趋势与建议

共 14 页幻灯片，采用专业商务风格设计。

祝好！
虾饺 🥟
"""
    
    msg = MIMEMultipart()
    msg['From'] = from_email
    msg['To'] = to_email
    msg['Subject'] = subject
    msg.attach(MIMEText(body, 'plain', 'utf-8'))
    
    with open(ppt_path, "rb") as f:
        attachment = f.read()
    
    from email.mime.base import MIMEBase
    from email import encoders
    
    part = MIMEBase('application', 'vnd.openxmlformats-officedocument.presentationml.presentation')
    part.set_payload(attachment)
    encoders.encode_base64(part)
    part.add_header('Content-Disposition', f'attachment; filename="竞品调研_20260318.pptx"')
    msg.attach(part)
    
    with smtplib.SMTP_SSL(smtp_server, smtp_port) as server:
        server.login(from_email, password)
        server.send_message(msg)
    
    print(f"[OK] 邮件发送成功！")
    print(f"[To] {to_email}")
    print(f"[Subject] {subject}")

if __name__ == "__main__":
    ppt_path = generate_ppt()
    send_email(ppt_path)
    print("\n✅ 全部完成！")

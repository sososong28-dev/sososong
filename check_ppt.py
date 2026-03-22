#!/usr/bin/env python3
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pathlib import Path

# 查找 PPTX 文件
ppt_files = list(Path('.').glob('*.pptx'))

if not ppt_files:
    print("未找到 PPTX 文件")
    sys.exit(1)

for ppt_file in ppt_files:
    print(f"\n文件：{ppt_file.name}")
    print(f"大小：{ppt_file.stat().st_size} 字节")
    
    try:
        prs = Presentation(ppt_file)
        print(f"格式：PPTX (可编辑)")
        print(f"幻灯片数：{len(prs.slides)}")
        print("\n页面列表:")
        for i, slide in enumerate(prs.slides, 1):
            title = slide.shapes.title.text if slide.shapes.title else "无标题"
            print(f"  {i}. {title}")
    except Exception as e:
        print(f"错误：{e}")
